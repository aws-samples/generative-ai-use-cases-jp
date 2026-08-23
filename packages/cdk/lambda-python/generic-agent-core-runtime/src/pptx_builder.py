"""PowerPoint rendering for the create_powerpoint tool.

Slides are described declaratively (a ``type`` plus its fields) and drawn onto
blank layouts, so the result does not depend on the template that python-pptx
ships with.

Two rules apply to every renderer:

* Geometry is derived from the space that is actually left on the slide, never
  from hard-coded heights. A key message pushes the body down, and the body has
  to absorb that.
* When the content cannot be drawn legibly, raise ``ValueError`` instead of
  letting shapes run off the slide. A caller can split the content and retry;
  it cannot notice silently clipped output.
"""

import math

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Palette. Two colours carry the design; the greys only separate information.
INK = RGBColor(0x23, 0x2F, 0x3E)
ACCENT = RGBColor(0xFF, 0x99, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xF7, 0xF8, 0xFA)
RULE = RGBColor(0xD5, 0xDB, 0xE1)
MUTED = RGBColor(0x6B, 0x74, 0x82)
ON_DARK = RGBColor(0xB9, 0xC3, 0xCF)

BODY_FONT = "Yu Gothic"
MONO_FONT = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.85)
FOOTER_TOP = SLIDE_H - Inches(0.62)
BODY_BOTTOM = FOOTER_TOP - Inches(0.1)

# Caps that keep a slide readable. Exceeding one is an error, not a resize.
MAX_MESSAGE_LINES = 2
MAX_SWIMLANE_COLUMNS = 10
MAX_SEQUENCE_MESSAGES = 10
MAX_SEQUENCE_ACTORS = 6
MAX_FLOW_SCALE_DOWN = 0.8
MAX_CODE_LINES = 18
# A tree going right can be deeper than one going down, which runs out of width.
MAX_TREE_DEPTH = {"right": 4, "down": 3}


def _set_font(run, size, bold=False, color=INK, name=BODY_FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    # A latin typeface alone leaves CJK glyphs on the theme font, so the East
    # Asian and complex-script slots have to be set as well.
    rpr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        el = rpr.find(A_NS + tag)
        if el is None:
            el = etree.SubElement(rpr, A_NS + tag)
        el.set("typeface", name)


def _rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_pt=1):
    # PowerPoint rejects a file that contains a shape with zero width or height,
    # which is easy to produce when a computed span collapses (a single child, an
    # empty range). Such a shape would be invisible anyway, so drop it.
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_pt)
    s.shadow.inherit = False
    if s.has_text_frame:
        s.text_frame.word_wrap = True
    return s


def _text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, spacing=1.2, anchor=MSO_ANCHOR.TOP):
    """Draw a text box. ``lines`` is a list of (text, size, bold, colour[, font])."""
    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, spec in enumerate(lines):
        body, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        font = spec[4] if len(spec) > 4 else BODY_FONT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = body
        _set_font(run, size, bold, color, font)
    return box


def _centered(slide, x, y, w, h, body, size, bold, color, font=BODY_FONT):
    """Vertically centre a single line inside a box of height ``h``."""
    line_h = Inches(size / 72 * 1.4)
    _text(slide, x, y + (h - line_h) / 2, w, line_h * 2, [(body, size, bold, color, font)], align=PP_ALIGN.CENTER, spacing=1.15)


def _h_arrow(slide, x1, x2, y, color=ACCENT, dashed=False, head=True):
    """Horizontal connector that may point either way."""
    lo, hi = (x1, x2) if x2 > x1 else (x2, x1)
    if dashed:
        seg, gap, x = Inches(0.07), Inches(0.05), lo
        while x < hi:
            _rect(slide, x, y - Pt(1), min(seg, hi - x), Pt(2), color)
            x += seg + gap
    else:
        _rect(slide, lo, y - Pt(1), hi - lo, Pt(2), color)
    if head:
        d = Inches(0.15)
        tip = _rect(slide, (x2 - d) if x2 > x1 else x2, y - d / 2, d, d, color, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
        tip.rotation = 90 if x2 > x1 else 270


def _v_arrow(slide, x, y1, y2, color=INK, head=True):
    _rect(slide, x - Pt(1), y1, Pt(2), y2 - y1, color)
    if head:
        d = Inches(0.16)
        _rect(slide, x - d / 2, y2 - d, d, d, color, shape=MSO_SHAPE.ISOSCELES_TRIANGLE).rotation = 180


def _lines_of(value):
    """Accept either a string with newlines or a list of strings."""
    if value is None:
        return []
    if isinstance(value, str):
        return [ln for ln in value.split("\n") if ln.strip()]
    return [str(v) for v in value if str(v).strip()]


def _declare_notes_master(prs):
    """Add the ``notesMasterIdLst`` that python-pptx leaves out.

    Creating a notes slide adds a notesMaster part and a relationship to it, but
    ``presentation.xml`` is never told about it. PowerPoint then refuses to open
    the file ("the file format is not valid"), so the declaration is written here.
    """
    rel_id = next((rid for rid, rel in prs.part.rels.items() if rel.reltype.endswith("/notesMaster")), None)
    if rel_id is None or prs._element.find(qn("p:notesMasterIdLst")) is not None:
        return
    lst = etree.SubElement(prs._element, qn("p:notesMasterIdLst"))
    entry = etree.SubElement(lst, qn("p:notesMasterId"))
    entry.set(qn("r:id"), rel_id)
    # The schema fixes the order of CT_Presentation's children.
    masters = prs._element.find(qn("p:sldMasterIdLst"))
    position = list(prs._element).index(masters) + 1 if masters is not None else 0
    prs._element.remove(lst)
    prs._element.insert(position, lst)


def _chrome(slide, spec, index, deck_title):
    """Draw kicker / title / accent rule / key message and the footer.

    Returns the y where the body may start. A key message tightens the header so
    the body keeps its height.
    """
    kicker = spec.get("kicker", "")
    title = spec.get("title", "")
    message = _lines_of(spec.get("message"))
    if len(message) > MAX_MESSAGE_LINES:
        raise ValueError(f"key message supports up to {MAX_MESSAGE_LINES} lines, got {len(message)}")

    tight = bool(message)
    kicker_y = Inches(0.55) if tight else Inches(0.62)
    title_y = Inches(0.88) if tight else Inches(0.98)
    title_pt = 25 if tight else 28
    rule_y = Inches(1.62) if tight else Inches(1.78)

    if kicker:
        _text(slide, MARGIN, kicker_y, SLIDE_W - MARGIN * 2, Inches(0.28), [(kicker.upper(), 10.5, True, ACCENT)])
    if title:
        _text(slide, MARGIN, title_y, SLIDE_W - MARGIN * 2, Inches(0.7), [(title, title_pt, True, INK)])
    _rect(slide, MARGIN, rule_y, Inches(0.95), Pt(3), ACCENT)

    if deck_title is not None:
        _rect(slide, MARGIN, FOOTER_TOP, SLIDE_W - MARGIN * 2, Pt(0.75), RULE)
        _text(slide, MARGIN, FOOTER_TOP + Inches(0.12), Inches(9), Inches(0.28), [(deck_title, 9, False, MUTED)])
        _text(slide, SLIDE_W - MARGIN - Inches(0.6), FOOTER_TOP + Inches(0.12), Inches(0.6), Inches(0.28), [(str(index), 9, False, MUTED)], align=PP_ALIGN.RIGHT)

    if not message:
        return Inches(2.05)

    message_y = rule_y + Inches(0.22)
    _text(slide, MARGIN, message_y, SLIDE_W - MARGIN * 2 - Inches(0.6), Inches(0.9), [(ln, 16, False, INK) for ln in message], spacing=1.35)
    return message_y + Inches(0.38) * len(message) + Inches(0.24)


def _fit(available, count, preferred, minimum, what):
    """Pitch that fits ``count`` items into ``available``, or an error."""
    if count <= 0:
        raise ValueError(f"{what} requires at least one item")
    pitch = min(preferred, available / count)
    if pitch < minimum:
        raise ValueError(f"{what} does not fit on one slide with {count} items; split it across slides")
    return pitch


# --------------------------------------------------------------------------
# Renderers. Each takes (slide, spec, top) where ``top`` is the body start y.
# --------------------------------------------------------------------------


def _render_title(slide, spec, top):
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, INK)
    _rect(slide, MARGIN, Inches(2.55), Inches(1.3), Pt(4), ACCENT)
    lines = _lines_of(spec.get("title")) or [""]
    _text(slide, MARGIN, Inches(2.95), SLIDE_W - MARGIN * 2, Inches(2.0), [(ln, 42, True, WHITE) for ln in lines], spacing=1.05)
    if spec.get("subtitle"):
        _text(slide, MARGIN, Inches(5.05), SLIDE_W - MARGIN * 2, Inches(0.5), [(spec["subtitle"], 17, False, ON_DARK)])
    if spec.get("footnote"):
        _text(slide, MARGIN, Inches(6.35), SLIDE_W - MARGIN * 2, Inches(0.35), [(spec["footnote"], 11, False, RGBColor(0x7E, 0x8C, 0x9E))])


def _render_section(slide, spec, top):
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, INK)
    if spec.get("number"):
        _text(slide, MARGIN, Inches(2.5), SLIDE_W - MARGIN * 2, Inches(1.2), [(str(spec["number"]), 64, True, ACCENT)])
    _rect(slide, MARGIN, Inches(3.95), Inches(1.3), Pt(4), ACCENT)
    _text(slide, MARGIN, Inches(4.35), SLIDE_W - MARGIN * 2, Inches(0.9), [(spec.get("title", ""), 34, True, WHITE)])


def _render_bullets(slide, spec, top):
    items = _lines_of(spec.get("bullets"))
    pitch = _fit(BODY_BOTTOM - top, len(items), Inches(0.62), Inches(0.34), "bullets")
    y = top
    for item in items:
        _rect(slide, MARGIN + Inches(0.02), y + Inches(0.13), Pt(5), Pt(5), ACCENT, shape=MSO_SHAPE.OVAL)
        _text(slide, MARGIN + Inches(0.34), y, SLIDE_W - MARGIN * 2 - Inches(0.4), pitch, [(item, 15, False, INK)], spacing=1.3)
        y += pitch


def _render_agenda(slide, spec, top):
    items = _lines_of(spec.get("items"))
    pitch = _fit(BODY_BOTTOM - top, len(items), Inches(0.88), Inches(0.46), "agenda")
    y = top
    for i, item in enumerate(items, 1):
        _text(slide, MARGIN, y, Inches(0.7), Inches(0.5), [(f"{i:02d}", 20, True, ACCENT)])
        _text(slide, MARGIN + Inches(0.95), y + Inches(0.06), Inches(9), Inches(0.5), [(item, 17, False, INK)])
        _rect(slide, MARGIN, y + pitch - Inches(0.26), SLIDE_W - MARGIN * 2, Pt(0.75), RULE)
        y += pitch


def _render_columns(slide, spec, top):
    cols = spec.get("columns") or []
    if not 2 <= len(cols) <= 4:
        raise ValueError(f"columns supports 2 to 4 columns, got {len(cols)}")
    gap = Inches(0.25)
    cw = (SLIDE_W - MARGIN * 2 - gap * (len(cols) - 1)) / len(cols)
    # Size the cards to their contents instead of a fixed height, so a short
    # list does not leave the bottom of the card empty.
    longest = max((len(_lines_of(c.get("items"))) for c in cols), default=0)
    row_pitch = _fit(BODY_BOTTOM - top - Inches(1.3), longest, Inches(0.66), Inches(0.34), "columns")
    height = min(Inches(1.0) + row_pitch * longest + Inches(0.3), BODY_BOTTOM - top)
    x = MARGIN
    for col in cols:
        rows = _lines_of(col.get("items"))
        highlight = bool(col.get("highlight"))
        _rect(slide, x, top, cw, height, WHITE, line=RULE)
        _rect(slide, x, top, cw, Pt(4), ACCENT if highlight else RULE)
        _text(slide, x + Inches(0.35), top + Inches(0.35), cw - Inches(0.7), Inches(0.4), [(col.get("title", ""), 16, True, INK)])
        y = top + Inches(1.0)
        for row in rows:
            _rect(slide, x + Inches(0.35), y + Inches(0.09), Pt(4.5), Pt(4.5), ACCENT, shape=MSO_SHAPE.OVAL)
            _text(slide, x + Inches(0.62), y, cw - Inches(1.0), row_pitch, [(row, 13, False, MUTED)], spacing=1.25)
            y += row_pitch
        x += cw + gap


def _render_kpi(slide, spec, top):
    items = spec.get("items") or []
    if not 2 <= len(items) <= 4:
        raise ValueError(f"kpi supports 2 to 4 metrics, got {len(items)}")
    gap = Inches(0.2)
    cw = (SLIDE_W - MARGIN * 2 - gap * (len(items) - 1)) / len(items)
    height = min(Inches(2.4), BODY_BOTTOM - top)
    x = MARGIN
    for item in items:
        _rect(slide, x, top, cw, height, WHITE, line=RULE)
        _rect(slide, x, top, cw, Pt(5), INK)
        _text(slide, x + Inches(0.35), top + Inches(0.4), cw - Inches(0.7), Inches(0.3), [(item.get("label", ""), 12, True, MUTED)])
        box = _text(slide, x + Inches(0.35), top + Inches(0.85), cw - Inches(0.7), Inches(0.9), [(str(item.get("value", "")), 44, True, INK)])
        if item.get("unit"):
            run = box.text_frame.paragraphs[0].add_run()
            run.text = "  " + str(item["unit"])
            _set_font(run, 15, False, MUTED)
        if item.get("note"):
            _text(slide, x + Inches(0.35), top + Inches(1.85), cw - Inches(0.7), Inches(0.4), [(item["note"], 11, False, ACCENT)])
        x += cw + gap


def _render_table(slide, spec, top):
    rows = spec.get("rows") or []
    if len(rows) < 2:
        raise ValueError("table requires a header row and at least one body row")
    cols = len(rows[0])
    if not all(len(r) == cols for r in rows):
        raise ValueError("every table row must have the same number of cells")
    pitch = _fit(BODY_BOTTOM - top, len(rows), Inches(0.58), Inches(0.34), "table")
    table = slide.shapes.add_table(len(rows), cols, int(MARGIN), int(top), int(SLIDE_W - MARGIN * 2), int(pitch)).table
    widths = spec.get("column_widths")
    if widths:
        if len(widths) != cols:
            raise ValueError("column_widths must have one entry per column")
        total = float(sum(widths))
        for c, w in enumerate(widths):
            table.columns[c].width = int((SLIDE_W - MARGIN * 2) * (w / total))
    for r, row in enumerate(rows):
        table.rows[r].height = int(pitch)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK if r == 0 else (WHITE if r % 2 else MIST)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.2)
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(value)
            _set_font(run, 11.5 if r == 0 else 11, r == 0, WHITE if r == 0 else INK)


def _render_chart(slide, spec, top):
    categories = _lines_of(spec.get("categories"))
    series = spec.get("series") or []
    if not categories or not series:
        raise ValueError("chart requires categories and at least one series")
    kinds = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED, "line": XL_CHART_TYPE.LINE_MARKERS, "pie": XL_CHART_TYPE.PIE}
    kind = spec.get("chart_type", "column")
    if kind not in kinds:
        raise ValueError(f"chart_type must be one of {sorted(kinds)}, got {kind!r}")

    insight = _lines_of(spec.get("insight"))
    height = min(Inches(4.2), BODY_BOTTOM - top)
    chart_w = SLIDE_W - MARGIN * 2 - (Inches(3.7) if insight else Inches(0))

    data = CategoryChartData()
    data.categories = categories
    for s in series:
        values = s.get("values") or []
        if len(values) != len(categories):
            raise ValueError("every series must have one value per category")
        data.add_series(s.get("name", ""), tuple(values))
    chart = slide.shapes.add_chart(kinds[kind], int(MARGIN), int(top), int(chart_w), int(height), data).chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    if kind != "pie":
        for i, plot_series in enumerate(chart.plots[0].series):
            plot_series.format.fill.solid()
            plot_series.format.fill.fore_color.rgb = INK if i == 0 else ACCENT

    if insight:
        ix = MARGIN + chart_w + Inches(0.3)
        iw = SLIDE_W - MARGIN - ix
        _rect(slide, ix, top, iw, height, MIST, line=RULE)
        _rect(slide, ix, top, Pt(4), height, ACCENT)
        # No default wording here: the deck's language belongs to the caller.
        body = []
        if spec.get("insight_title"):
            body.append((spec["insight_title"], 13, True, INK))
        body += [(ln, 12, False, MUTED, BODY_FONT) for ln in insight]
        _text(slide, ix + Inches(0.3), top + Inches(0.35), iw - Inches(0.6), height - Inches(0.7), body, spacing=1.35)


def _render_architecture(slide, spec, top):
    nodes = spec.get("nodes") or []
    if not 2 <= len(nodes) <= 5:
        raise ValueError(f"architecture supports 2 to 5 nodes, got {len(nodes)}")
    gap = Inches(0.6)
    nw = (SLIDE_W - MARGIN * 2 - gap * (len(nodes) - 1)) / len(nodes)
    nh = Inches(1.55)
    y = top + Inches(0.5)
    for i, node in enumerate(nodes):
        x = MARGIN + (nw + gap) * i
        hot = bool(node.get("highlight"))
        _rect(slide, x, y, nw, nh, INK if hot else MIST, line=None if hot else RULE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(slide, x, y + Inches(0.3), nw, Inches(0.85), [(ln, 15, True, WHITE if hot else INK) for ln in _lines_of(node.get("label"))], align=PP_ALIGN.CENTER, spacing=1.15)
        if node.get("note"):
            _text(slide, x, y + nh + Inches(0.16), nw, Inches(0.35), [(node["note"], 11, False, MUTED)], align=PP_ALIGN.CENTER)
        if i:
            ax = x - gap
            _h_arrow(slide, ax + Inches(0.12), x - Inches(0.05), y + nh / 2)


def _render_code(slide, spec, top):
    code = _lines_of(spec.get("code")) if not isinstance(spec.get("code"), str) else spec["code"].split("\n")
    if len(code) > MAX_CODE_LINES:
        raise ValueError(f"code supports up to {MAX_CODE_LINES} lines, got {len(code)}; show an excerpt")
    caption = spec.get("caption")
    height = min(Inches(3.85), BODY_BOTTOM - top - (Inches(0.45) if caption else Inches(0)))
    _rect(slide, MARGIN, top, SLIDE_W - MARGIN * 2, height, INK)
    _rect(slide, MARGIN, top, Pt(4), height, ACCENT)
    box = slide.shapes.add_textbox(int(MARGIN + Inches(0.45)), int(top + Inches(0.3)), int(SLIDE_W - MARGIN * 2 - Inches(0.9)), int(height - Inches(0.6)))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    comment_prefixes = ("#", "//", "--")
    for i, line in enumerate(code):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.35
        num = p.add_run()
        num.text = f"{i + 1:>2}  "
        _set_font(num, 10, False, RGBColor(0x5A, 0x6B, 0x7E), MONO_FONT)
        body = p.add_run()
        body.text = line
        is_comment = line.strip().startswith(comment_prefixes)
        _set_font(body, 11.5, False, ACCENT if is_comment else RGBColor(0xE8, 0xEE, 0xF4), MONO_FONT)
    if caption:
        _text(slide, MARGIN, top + height + Inches(0.16), SLIDE_W - MARGIN * 2, Inches(0.4), [(caption, 11, False, MUTED)])


def _render_swimlane(slide, spec, top):
    """Roadmap swimlane: rows are workstreams, columns are periods."""
    columns = _lines_of(spec.get("columns"))
    lanes = spec.get("lanes") or []
    if len(columns) > MAX_SWIMLANE_COLUMNS:
        raise ValueError(f"swimlane supports up to {MAX_SWIMLANE_COLUMNS} columns, got {len(columns)}")
    if not columns or not lanes:
        raise ValueError("swimlane requires columns and lanes")

    milestones = spec.get("milestones") or []
    label_w = Inches(2.15)
    track_x = MARGIN + label_w
    track_w = SLIDE_W - MARGIN * 2 - label_w
    col_w = track_w / len(columns)
    header_h = Inches(0.42)
    lanes_top = top + header_h
    milestone_h = Inches(0.78) if milestones else Inches(0)
    lane_h = _fit(BODY_BOTTOM - lanes_top - milestone_h, len(lanes), Inches(0.86), Inches(0.44), "swimlane")

    _rect(slide, track_x, top, track_w, header_h, INK)
    for i, name in enumerate(columns):
        _centered(slide, track_x + col_w * i, top, col_w, header_h, name, 11, True, WHITE)
        if i:
            _rect(slide, track_x + col_w * i, top, Pt(0.75), header_h, RGBColor(0x3C, 0x4A, 0x5B))

    for r, lane in enumerate(lanes):
        y = lanes_top + lane_h * r
        shade = MIST if r % 2 == 0 else WHITE
        _rect(slide, MARGIN, y, label_w, lane_h, shade, line=RULE)
        _text(slide, MARGIN + Inches(0.28), y + (lane_h - Inches(0.3)) / 2, label_w - Inches(0.5), Inches(0.35), [(lane.get("name", ""), 13, True, INK)])
        _rect(slide, track_x, y, track_w, lane_h, shade, line=RULE)
        for i in range(1, len(columns)):
            _rect(slide, track_x + col_w * i, y, Pt(0.75), lane_h, RULE)
        for bar in lane.get("bars") or []:
            start, span = int(bar.get("start", 0)), int(bar.get("span", 1))
            if start < 0 or start + span > len(columns):
                raise ValueError(f"bar {bar} falls outside the {len(columns)} columns")
            tone = {"accent": ACCENT, "muted": MUTED}.get(bar.get("tone"), INK)
            bx = track_x + col_w * start + Inches(0.06)
            bw = col_w * span - Inches(0.12)
            bh = lane_h * 0.5
            by = y + (lane_h - bh) / 2
            _rect(slide, bx, by, bw, bh, tone, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _centered(slide, bx, by, bw, bh, bar.get("label", ""), 11, True, WHITE)

    bottom = lanes_top + lane_h * len(lanes)
    for ms in milestones:
        col = int(ms.get("column", 0))
        mx = track_x + col_w * col + col_w / 2
        _rect(slide, mx - Inches(0.11), bottom + Inches(0.1), Inches(0.22), Inches(0.22), ACCENT, shape=MSO_SHAPE.DIAMOND)
        _text(slide, mx - Inches(1.0), bottom + Inches(0.4), Inches(2.0), Inches(0.3), [(ms.get("label", ""), 10, True, ACCENT)], align=PP_ALIGN.CENTER)
    if spec.get("today") is not None:
        nx = track_x + col_w * float(spec["today"])
        _rect(slide, nx, top, Pt(2), bottom - top, ACCENT)
        if spec.get("today_label"):
            _text(slide, nx - Inches(0.6), top - Inches(0.34), Inches(1.2), Inches(0.3), [(spec["today_label"], 10, True, ACCENT)], align=PP_ALIGN.CENTER)


def _render_sequence(slide, spec, top):
    actors = _lines_of(spec.get("actors"))
    messages = spec.get("messages") or []
    if not 2 <= len(actors) <= MAX_SEQUENCE_ACTORS:
        raise ValueError(f"sequence supports 2 to {MAX_SEQUENCE_ACTORS} actors, got {len(actors)}")
    if len(messages) > MAX_SEQUENCE_MESSAGES:
        raise ValueError(f"sequence supports up to {MAX_SEQUENCE_MESSAGES} messages, got {len(messages)}")

    box_w, box_h = Inches(1.95), Inches(0.62)
    gap = (SLIDE_W - MARGIN * 2 - box_w * len(actors)) / (len(actors) - 1)
    highlight = spec.get("highlight")
    centres = []
    for i, actor in enumerate(actors):
        x = MARGIN + (box_w + gap) * i
        hot = actor == highlight
        _rect(slide, x, top, box_w, box_h, INK if hot else MIST, line=None if hot else RULE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _centered(slide, x, top, box_w, box_h, actor, 12, True, WHITE if hot else INK)
        centres.append(x + box_w / 2)

    life_top = top + box_h
    first_y = life_top + Inches(0.45)
    step = min(Inches(0.62), (BODY_BOTTOM - Inches(0.28) - first_y) / max(1, len(messages) - 1))
    if step < Inches(0.34):
        raise ValueError("sequence does not fit on one slide; split the messages")
    life_bottom = first_y + step * (len(messages) - 1) + Inches(0.3)

    for x in centres:
        y = life_top
        while y < life_bottom:
            _rect(slide, x - Pt(0.75), y, Pt(1.5), Inches(0.07), RULE)
            y += Inches(0.13)

    if spec.get("legend"):
        _text(slide, SLIDE_W - MARGIN - Inches(4.5), Inches(1.72), Inches(4.5), Inches(0.3), [(spec["legend"], 10.5, False, MUTED)], align=PP_ALIGN.RIGHT)

    y = first_y
    for msg in messages:
        a, b = int(msg.get("from", 0)), int(msg.get("to", 0))
        if not (0 <= a < len(actors) and 0 <= b < len(actors)):
            raise ValueError(f"message {msg} refers to an actor index that does not exist")
        label, ret = msg.get("label", ""), bool(msg.get("return"))
        if a == b:
            x = centres[a]
            w, h = Inches(0.5), Inches(0.26)
            _rect(slide, x, y - h / 2, w, Pt(2), ACCENT)
            _rect(slide, x + w, y - h / 2, Pt(2), h, ACCENT)
            _h_arrow(slide, x + w, x, y + h / 2, ACCENT)
            # The label goes to the right of the loop, or to its left when the
            # rightmost actor leaves no room.
            label_x = x + w + Inches(0.16)
            room = SLIDE_W - MARGIN - label_x
            if room >= Inches(1.2):
                _text(slide, label_x, y - Inches(0.1), room, Inches(0.3), [(label, 10, False, MUTED)])
            else:
                width = Inches(2.6)
                _text(slide, x - Inches(0.16) - width, y - Inches(0.1), width, Inches(0.3), [(label, 10, False, MUTED)], align=PP_ALIGN.RIGHT)
        else:
            _h_arrow(slide, centres[a], centres[b], y, MUTED if ret else ACCENT, dashed=ret)
            lo, hi = sorted((centres[a], centres[b]))
            _text(slide, lo, y - Inches(0.33), hi - lo, Inches(0.28), [(label, 10, not ret, MUTED if ret else INK)], align=PP_ALIGN.CENTER)
        y += step


def _tree_depth(node):
    kids = node.get("children") or []
    return 1 + max((_tree_depth(k) for k in kids), default=0)


def _render_logictree(slide, spec, top):
    """A tree of labelled nodes.

    ``direction`` decides what it reads as: "right" decomposes a question from
    the left (an issue tree), "down" hangs children under a parent (an org
    chart). Both share the layout rule that a subtree claims the space its
    leaves need, so an uneven tree still lines up.
    """
    root = spec.get("root")
    if not root:
        raise ValueError("logictree requires a root node")
    direction = spec.get("direction", "right")
    if direction not in ("right", "down"):
        raise ValueError(f'direction must be "right" or "down", got {direction!r}')
    depth = _tree_depth(root)
    max_depth = MAX_TREE_DEPTH[direction]
    if depth > max_depth:
        raise ValueError(f"logictree going {direction} supports up to {max_depth} levels, got {depth}")

    def leaves(node):
        kids = node.get("children") or []
        return sum(leaves(k) for k in kids) if kids else 1

    available = BODY_BOTTOM - top - Inches(0.1)
    leaf_count = leaves(root)

    if direction == "right":
        gap = Inches(0.55)
        level_w = (SLIDE_W - MARGIN * 2 - gap * (depth - 1)) / depth
        band = available / leaf_count
        if band < Inches(0.52):
            raise ValueError(f"logictree has {leaf_count} leaves and does not fit; split the branches")
        node_h = min(Inches(0.66), band - Inches(0.12))
        node_w = level_w
    else:
        gap = Inches(0.28)
        node_w = Inches(2.05)
        span = node_w * leaf_count + gap * (leaf_count - 1)
        if span > SLIDE_W - MARGIN * 2:
            raise ValueError(f"logictree has {leaf_count} leaves and is too wide; split the branches")
        node_h = min(Inches(0.82), available / depth * 0.58)
        level_gap = (available - node_h * depth) / max(1, depth - 1)
        levels = [top + Inches(0.15) + (node_h + level_gap) * d for d in range(depth)]

    def extent(node):
        """How much of the cross axis this subtree needs."""
        kids = node.get("children") or []
        if not kids:
            return node_w if direction == "down" else available / leaf_count
        return sum(extent(k) for k in kids) + (gap if direction == "down" else 0) * (len(kids) - 1)

    def draw(node, offset, size, level):
        kids = node.get("children") or []
        is_root = level == 0
        if direction == "right":
            x = MARGIN + (node_w + gap) * level
            centre = offset + size / 2
            y = centre - node_h / 2
        else:
            x = offset + (size - node_w) / 2
            y = levels[level]
            centre = x + node_w / 2

        _rect(slide, x, y, node_w, node_h, INK if is_root else WHITE, line=None if is_root else RULE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        accent = ACCENT if (node.get("highlight") or (direction == "down" and level == 1)) else RULE
        if not is_root:
            if direction == "right":
                _rect(slide, x, y, Pt(4), node_h, accent)
            else:
                _rect(slide, x, y, node_w, Pt(3), accent)

        label = node.get("label", "")
        note = node.get("note")
        bold = is_root or bool(node.get("highlight"))
        colour = WHITE if is_root else INK
        if direction == "down":
            _text(slide, x, y + Inches(0.15), node_w, Inches(0.3), [(label, 12, True, colour)], align=PP_ALIGN.CENTER)
            if note:
                _text(slide, x, y + Inches(0.46), node_w, Inches(0.26), [(note, 9.5, False, ON_DARK if is_root else MUTED)], align=PP_ALIGN.CENTER)
        else:
            _text(slide, x + Inches(0.22), y + (node_h - Inches(0.26)) / 2, node_w - Inches(0.44), Inches(0.5), [(label, 11.5, bold, colour)], spacing=1.2)

        if not kids:
            return
        # Elbow: out of the parent, along a shared spine, into each child.
        centres, cursor = [], offset
        child_specs = []
        for kid in kids:
            share = extent(kid) if direction == "down" else size * leaves(kid) / leaves(node)
            child_specs.append((kid, cursor, share))
            cursor += share + (gap if direction == "down" else 0)
        if direction == "right":
            spine = x + node_w + gap / 2
            _rect(slide, x + node_w, centre - Pt(1), spine - (x + node_w), Pt(2), RULE)
            for _, kid_offset, share in child_specs:
                centres.append(kid_offset + share / 2)
            if len(centres) > 1:
                _rect(slide, spine - Pt(1), min(centres), Pt(2), max(centres) - min(centres), RULE)
            child_x = MARGIN + (node_w + gap) * (level + 1)
            for c in centres:
                _h_arrow(slide, spine, child_x - Inches(0.02), c, RULE, head=False)
        else:
            spine_y = y + node_h + level_gap / 2
            _rect(slide, centre - Pt(1), y + node_h, Pt(2), spine_y - (y + node_h), RULE)
            for _, kid_offset, share in child_specs:
                centres.append(kid_offset + (share - node_w) / 2 + node_w / 2)
            if len(centres) > 1:
                _rect(slide, min(centres), spine_y - Pt(1), max(centres) - min(centres), Pt(2), RULE)
            for c in centres:
                _rect(slide, c - Pt(1), spine_y, Pt(2), levels[level + 1] - spine_y, RULE)
        for kid, kid_offset, share in child_specs:
            draw(kid, kid_offset, share, level + 1)

    if direction == "right":
        draw(root, top + Inches(0.05), available, 0)
    else:
        span = extent(root)
        draw(root, MARGIN + (SLIDE_W - MARGIN * 2 - span) / 2, span, 0)


def _render_flow(slide, spec, top):
    """Decision flow running down the slide.

    A long flow wraps into a second column rather than shrinking until it is
    unreadable; ``columns`` pins the choice when the caller wants one shape.
    """
    steps = spec.get("steps") or []
    if not steps:
        raise ValueError("flow requires at least one step")
    requested = spec.get("columns")
    if requested not in (None, 1, 2):
        raise ValueError(f"flow supports 1 or 2 columns, got {requested!r}")

    base_h, base_diamond, base_gap = Inches(0.66), Inches(0.92), Inches(0.34)
    available = BODY_BOTTOM - Inches(0.06) - top

    def height_of(step):
        return base_diamond if step.get("type") == "decision" else base_h

    def split(count):
        per = math.ceil(len(steps) / count)
        chunks = [steps[i * per : (i + 1) * per] for i in range(count)]
        return [c for c in chunks if c]

    def scale_for(chunks):
        tallest = max(sum(height_of(s) for s in chunk) + base_gap * (len(chunk) - 1) for chunk in chunks)
        return min(1.0, available / tallest)

    for count in [requested] if requested else (1, 2):
        chunks = split(count)
        scale = scale_for(chunks)
        if scale >= MAX_FLOW_SCALE_DOWN:
            break
    else:
        raise ValueError(f"flow does not fit on one slide ({len(steps)} steps); split it across slides")
    if len(chunks) == 1:
        node_w, lane_w = Inches(2.9), SLIDE_W - MARGIN * 2
        side_w, side_offset = Inches(2.7), Inches(4.1)
        first_cx = MARGIN + Inches(3.6)
    else:
        lane_w = (SLIDE_W - MARGIN * 2) / 2
        node_w, side_w = Inches(2.5), Inches(2.2)
        side_offset = Inches(2.9)
        first_cx = MARGIN + Inches(1.45)

    nh, dh, gap = base_h * scale, base_diamond * scale, base_gap * scale
    column_bounds = []
    for index, chunk in enumerate(chunks):
        cx = first_cx + lane_w * index
        side_cx = cx + side_offset
        y = top
        previous_bottom = None
        first_kind = chunk[0].get("type", "process")
        entry_y = top + (dh if first_kind == "decision" else nh) / 2
        for step in chunk:
            kind = step.get("type", "process")
            h = dh if kind == "decision" else nh
            x = cx - node_w / 2
            if previous_bottom is not None:
                _v_arrow(slide, cx, previous_bottom, y)
            if kind == "decision":
                _rect(slide, x, y, node_w, h, WHITE, line=ACCENT, shape=MSO_SHAPE.DIAMOND, line_pt=1.75)
                _centered(slide, x + Inches(0.2), y, node_w - Inches(0.4), h, step.get("label", ""), 10.5, True, INK)
                by = y + h / 2
                _h_arrow(slide, cx + node_w / 2, side_cx - side_w / 2, by, ACCENT)
                _text(slide, cx + node_w / 2 + Inches(0.1), by - Inches(0.32), Inches(1.1), Inches(0.28), [(step.get("no_label", "No"), 9.5, True, ACCENT)])
                _rect(slide, side_cx - side_w / 2, by - Inches(0.29), side_w, Inches(0.58), MIST, line=RULE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                _centered(slide, side_cx - side_w / 2, by - Inches(0.29), side_w, Inches(0.58), step.get("no_branch", ""), 10, False, MUTED)
                _text(slide, cx + Inches(0.12), y + h - Inches(0.04), Inches(1.2), Inches(0.26), [(step.get("yes_label", "Yes"), 9.5, True, INK)])
            else:
                terminal = kind in ("start", "end")
                _rect(slide, x, y, node_w, h, INK if terminal else WHITE, line=None if terminal else RULE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                if not terminal:
                    _rect(slide, x, y, node_w, Pt(3), ACCENT)
                _centered(slide, x, y, node_w, h, step.get("label", ""), 11.5, True, WHITE if terminal else INK)
            previous_bottom = y + h
            y = previous_bottom + gap
        column_bounds.append((cx, previous_bottom, entry_y))

    # Carry the flow from the foot of one column into the side of the next one's
    # first node. Routing over the top would cross the column it is entering.
    for index in range(len(column_bounds) - 1):
        from_cx, from_bottom, _ = column_bounds[index]
        to_cx, _, to_entry = column_bounds[index + 1]
        to_edge = to_cx - node_w / 2
        gutter = to_edge - Inches(0.3)
        elbow = from_bottom + Inches(0.28)
        _rect(slide, from_cx - Pt(1), from_bottom, Pt(2), elbow - from_bottom, INK)
        _rect(slide, from_cx, elbow - Pt(1), gutter - from_cx, Pt(2), INK)
        _rect(slide, gutter - Pt(1), to_entry, Pt(2), elbow - to_entry, INK)
        _h_arrow(slide, gutter, to_edge, to_entry, INK)


def _render_quadrant(slide, spec, top):
    axis_h = Inches(0.55)
    fx = MARGIN + Inches(1.05)
    fw = Inches(9.1)
    fh = min(Inches(4.15), BODY_BOTTOM - axis_h - top - Inches(0.15))
    if fh < Inches(2.2):
        raise ValueError("quadrant does not fit; shorten the key message")
    fy = top + Inches(0.15)
    cx, cy = fx + fw / 2, fy + fh / 2

    tints = [(0, 0, MIST), (1, 0, RGBColor(0xFF, 0xF6, 0xE8)), (0, 1, RGBColor(0xFA, 0xFB, 0xFC)), (1, 1, MIST)]
    for qx, qy, tint in tints:
        _rect(slide, fx + fw / 2 * qx, fy + fh / 2 * qy, fw / 2, fh / 2, tint)
    _rect(slide, fx, fy, fw, fh, None, line=RULE)
    _rect(slide, cx - Pt(0.75), fy, Pt(1.5), fh, RULE)
    _rect(slide, fx, cy - Pt(0.75), fw, Pt(1.5), RULE)

    # Quadrant captions sit in the outer corners so they never touch the axes.
    captions = spec.get("quadrants") or {}
    corners = {"top_left": (fx + Inches(0.22), fy + Inches(0.2)), "top_right": (cx + Inches(0.22), fy + Inches(0.2)), "bottom_left": (fx + Inches(0.22), fy + fh - Inches(0.62)), "bottom_right": (cx + Inches(0.22), fy + fh - Inches(0.62))}
    for key, (lx, ly) in corners.items():
        lines = _lines_of(captions.get(key))
        if not lines:
            continue
        _text(slide, lx, ly - Inches(0.2) * (len(lines) - 1), Inches(3.6), Inches(0.8), [(ln, 11.5, True, MUTED) for ln in lines], spacing=1.25)

    if spec.get("x_axis"):
        _text(slide, fx, fy + fh + Inches(0.22), fw, Inches(0.3), [(spec["x_axis"], 11, True, INK)], align=PP_ALIGN.CENTER)
    if spec.get("y_axis"):
        aw = Inches(3.0)
        box = slide.shapes.add_textbox(int(fx - Inches(0.52) - aw / 2), int(cy - Inches(0.18)), int(aw), int(Inches(0.36)))
        box.rotation = 270
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = spec["y_axis"]
        _set_font(run, 11, True, INK)

    for point in spec.get("points") or []:
        u, v = float(point.get("x", 0.5)), float(point.get("y", 0.5))
        if not (0 <= u <= 1 and 0 <= v <= 1):
            raise ValueError(f"point {point} must have x and y between 0 and 1")
        px, py = fx + fw * u, fy + fh * (1 - v)
        hot = bool(point.get("highlight"))
        d = Inches(0.2) if hot else Inches(0.15)
        _rect(slide, px - d / 2, py - d / 2, d, d, ACCENT if hot else INK, shape=MSO_SHAPE.OVAL)
        flip = u > 0.68
        tw = Inches(2.6)
        _text(slide, px - tw - Inches(0.2) if flip else px + Inches(0.2), py - Inches(0.13), tw, Inches(0.3), [(point.get("label", ""), 11.5, hot, INK)], align=PP_ALIGN.RIGHT if flip else PP_ALIGN.LEFT)


# Slides drawn edge to edge; they carry no heading, footer or page number.
FULL_BLEED = {"title", "section"}

RENDERERS = {
    "title": _render_title,
    "section": _render_section,
    "agenda": _render_agenda,
    "bullets": _render_bullets,
    "kpi": _render_kpi,
    "chart": _render_chart,
    "table": _render_table,
    "columns": _render_columns,
    "quadrant": _render_quadrant,
    "architecture": _render_architecture,
    "sequence": _render_sequence,
    "swimlane": _render_swimlane,
    "flow": _render_flow,
    "logictree": _render_logictree,
    "code": _render_code,
}

SLIDE_TYPES = sorted(RENDERERS)


def build_presentation(slides, deck_title=None):
    """Render ``slides`` and return the Presentation.

    Each entry is a dict with a ``type`` from :data:`SLIDE_TYPES` plus the fields
    that type needs. ``title``, ``kicker`` and ``message`` are common to every
    non full-bleed type.
    """
    if not slides:
        raise ValueError("at least one slide is required")

    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)
    blank = prs.slide_layouts[6]

    page = 0
    for i, spec in enumerate(slides):
        if not isinstance(spec, dict):
            raise ValueError(f"slide {i + 1} must be an object, got {type(spec).__name__}")
        kind = spec.get("type", "bullets")
        render = RENDERERS.get(kind)
        if render is None:
            raise ValueError(f"slide {i + 1}: unknown type {kind!r}; expected one of {SLIDE_TYPES}")
        slide = prs.slides.add_slide(blank)
        try:
            if kind in FULL_BLEED:
                render(slide, spec, Inches(2.05))
            else:
                page += 1
                top = _chrome(slide, spec, page, deck_title)
                render(slide, spec, top)
        except ValueError as e:
            raise ValueError(f"slide {i + 1} ({kind}): {e}") from e

        notes = spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    _declare_notes_master(prs)
    return prs

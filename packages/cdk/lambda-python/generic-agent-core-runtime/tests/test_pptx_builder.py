"""Tests for the PowerPoint builder.

The costly failure mode here is a deck that renders without raising but that
PowerPoint refuses to open, so most of these assertions inspect the saved
package rather than the in-memory objects.
"""

import re
import zipfile

import pytest
from lxml import etree
from pptx.util import Inches

from src.pptx_builder import SLIDE_TYPES, build_presentation

P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER_TOP = SLIDE_H - Inches(0.62)
TOLERANCE = Inches(0.03)

# One minimal but valid spec per slide type, so the suite fails when a new type
# is added without coverage.
SPECS = {
    "title": {"type": "title", "title": "表題\n2行目", "subtitle": "副題", "footnote": "脚注"},
    "section": {"type": "section", "number": "02", "title": "章"},
    "bullets": {"type": "bullets", "title": "見出し", "bullets": ["あ", "い", "う"]},
    "agenda": {"type": "agenda", "title": "見出し", "items": ["一", "二", "三"]},
    "columns": {
        "type": "columns",
        "title": "見出し",
        "columns": [
            {"title": "左", "items": ["a", "b"]},
            {"title": "右", "items": ["c"], "highlight": True},
        ],
    },
    "kpi": {
        "type": "kpi",
        "title": "見出し",
        "items": [
            {"label": "売上", "value": "100", "unit": "千円", "note": "前年比"},
            {"label": "粗利", "value": "28", "unit": "%"},
        ],
    },
    "table": {
        "type": "table",
        "title": "見出し",
        "rows": [["列1", "列2"], ["値1", "値2"], ["値3", "値4"]],
        "column_widths": [2, 1],
    },
    "chart": {
        "type": "chart",
        "title": "見出し",
        "categories": ["1月", "2月"],
        "series": [{"name": "売上", "values": [1, 2]}],
        "insight": ["伸びている"],
        "insight_title": "読み取れること",
    },
    "architecture": {
        "type": "architecture",
        "title": "見出し",
        "nodes": [{"label": "A", "note": "説明"}, {"label": "B", "highlight": True}],
    },
    "code": {"type": "code", "title": "見出し", "code": "x = 1\n# コメント", "caption": "説明"},
    "swimlane": {
        "type": "swimlane",
        "title": "見出し",
        "columns": ["8月", "9月", "10月"],
        "lanes": [
            {"name": "設計", "bars": [{"start": 0, "span": 2, "label": "作業"}]},
            {"name": "実装", "bars": [{"start": 1, "span": 1, "label": "作業", "tone": "accent"}]},
        ],
        "milestones": [{"column": 2, "label": "公開"}],
        "today": 1.5,
        "today_label": "現在",
    },
    "sequence": {
        "type": "sequence",
        "title": "見出し",
        "actors": ["利用者", "UI", "Runtime"],
        "highlight": "Runtime",
        "legend": "実線が呼び出し、点線が戻り",
        "messages": [
            {"from": 0, "to": 1, "label": "依頼"},
            {"from": 1, "to": 2, "label": "呼び出し"},
            {"from": 2, "to": 2, "label": "自己処理"},
            {"from": 2, "to": 0, "label": "返却", "return": True},
        ],
    },
    "logictree": {
        "type": "logictree",
        "title": "見出し",
        "direction": "right",
        "root": {
            "label": "問い",
            "children": [
                {"label": "枝1", "highlight": True, "children": [{"label": "葉1"}]},
                {"label": "枝2", "children": [{"label": "葉2"}, {"label": "葉3"}]},
            ],
        },
    },
    "logictree_down": {
        "type": "logictree",
        "direction": "down",
        "title": "見出し",
        "root": {
            "label": "親",
            "note": "説明",
            # A node with one child used to collapse the connector to zero size.
            "children": [
                {"label": "子1", "children": [{"label": "孫1"}, {"label": "孫2"}]},
                {"label": "子2", "children": [{"label": "孫3"}]},
            ],
        },
    },
    "flow": {
        "type": "flow",
        "title": "見出し",
        "steps": [
            {"type": "start", "label": "開始"},
            {"type": "decision", "label": "条件", "no_branch": "別処理"},
            {"type": "end", "label": "終了"},
        ],
    },
    "quadrant": {
        "type": "quadrant",
        "title": "見出し",
        "x_axis": "横軸",
        "y_axis": "縦軸",
        "quadrants": {"top_left": "左上\n2行", "top_right": "右上"},
        "points": [
            {"label": "点1", "x": 0.2, "y": 0.3},
            {"label": "点2", "x": 0.8, "y": 0.7, "highlight": True},
        ],
    },
}


@pytest.fixture(scope="module")
def all_specs():
    return dict(SPECS)


def save(prs, tmp_path, name="deck.pptx"):
    path = tmp_path / name
    prs.save(str(path))
    return path


def parts(path):
    z = zipfile.ZipFile(path)
    return z, set(z.namelist())


def test_every_slide_type_has_a_spec(all_specs):
    """`logictree_down` covers the same renderer through its other direction."""
    covered = {spec["type"] for spec in all_specs.values()}
    assert sorted(covered) == SLIDE_TYPES


@pytest.mark.parametrize("kind", sorted(SPECS))
def test_each_type_renders_without_degenerate_shapes(kind, tmp_path):
    """A shape with zero width or height makes PowerPoint reject the file."""
    prs = build_presentation([SPECS[kind]], "テスト")
    path = save(prs, tmp_path, f"{kind}.pptx")
    z, names = parts(path)
    for name in names:
        if re.match(r"ppt/slides/slide\d+\.xml$", name):
            root = etree.fromstring(z.read(name))
            for ext in root.iter(A_NS + "ext"):
                assert ext.get("cx") != "0" and ext.get("cy") != "0", f"{kind}: zero-sized shape"


@pytest.mark.parametrize("kind", sorted(SPECS))
def test_each_type_stays_inside_the_slide(kind, tmp_path):
    prs = build_presentation([SPECS[kind]], "テスト")
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.rotation:
            continue  # a rotated box is placed by its centre
        assert shape.left >= -TOLERANCE, f"{kind}: shape starts left of the slide"
        assert shape.left + shape.width <= SLIDE_W + TOLERANCE, f"{kind}: shape runs off the right"
        assert shape.top + shape.height <= SLIDE_H + TOLERANCE, f"{kind}: shape runs off the bottom"


@pytest.mark.parametrize("kind", sorted(SPECS))
def test_body_never_reaches_the_footer(kind, tmp_path):
    """Footer furniture is 9pt; anything larger must stop above the rule."""
    prs = build_presentation([SPECS[kind]], "テスト")
    for shape in prs.slides[0].shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
        if runs and runs[0].font.size and runs[0].font.size.pt <= 9.5:
            continue
        assert shape.top + shape.height <= FOOTER_TOP + TOLERANCE, f"{kind}: {shape.text_frame.text[:20]!r} overlaps the footer"


def test_all_parts_are_well_formed_xml(all_specs, tmp_path):
    prs = build_presentation(list(all_specs.values()), "テスト")
    z, names = parts(save(prs, tmp_path))
    for name in names:
        if name.endswith((".xml", ".rels")):
            etree.fromstring(z.read(name))


def test_notes_declare_a_notes_master(tmp_path):
    """python-pptx adds the notesMaster relationship but not the declaration,
    and PowerPoint then reports the file as invalid."""
    prs = build_presentation([{"type": "bullets", "title": "見出し", "bullets": ["あ"], "notes": "メモ"}])
    z, names = parts(save(prs, tmp_path))
    assert any(n.startswith("ppt/notesSlides/notesSlide") for n in names)
    presentation = etree.fromstring(z.read("ppt/presentation.xml"))
    declared = presentation.find(P_NS + "notesMasterIdLst")
    assert declared is not None, "notesMasterIdLst is missing"
    children = [etree.QName(c).localname for c in presentation]
    # The schema fixes the order of CT_Presentation's children.
    assert children.index("notesMasterIdLst") == children.index("sldMasterIdLst") + 1
    assert children.index("notesMasterIdLst") < children.index("sldIdLst")


def test_no_notes_means_no_notes_master(tmp_path):
    prs = build_presentation([{"type": "bullets", "title": "見出し", "bullets": ["あ"]}])
    z, _ = parts(save(prs, tmp_path))
    presentation = etree.fromstring(z.read("ppt/presentation.xml"))
    assert presentation.find(P_NS + "notesMasterIdLst") is None


def test_notes_text_is_stored(tmp_path):
    prs = build_presentation([{"type": "bullets", "title": "見出し", "bullets": ["あ"], "notes": "話す内容"}])
    assert prs.slides[0].notes_slide.notes_text_frame.text == "話す内容"


def test_newlines_become_separate_paragraphs(tmp_path):
    prs = build_presentation([{"type": "title", "title": "一行目\n二行目"}])
    texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
    assert "一行目\n二行目" in texts
    assert not any("\\n" in t for t in texts), "a literal backslash-n reached the slide"


def test_key_message_pushes_the_body_down(tmp_path):
    without = build_presentation([{"type": "bullets", "title": "見出し", "bullets": ["あ"]}])
    with_message = build_presentation([{"type": "bullets", "title": "見出し", "message": "一行目\n二行目", "bullets": ["あ"]}])

    def bullet_top(prs):
        return next(s.top for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text == "あ")

    assert bullet_top(with_message) > bullet_top(without)


def test_key_message_is_larger_than_the_body(tmp_path):
    prs = build_presentation([{"type": "bullets", "title": "見出し", "message": "主張", "bullets": ["あ"]}])
    sizes = {s.text_frame.text: s.text_frame.paragraphs[0].runs[0].font.size.pt for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.paragraphs[0].runs}
    assert sizes["主張"] > sizes["あ"]


def test_deck_title_and_page_numbers_appear_only_on_body_slides(tmp_path):
    prs = build_presentation(
        [
            {"type": "title", "title": "表紙"},
            {"type": "bullets", "title": "本文1", "bullets": ["あ"]},
            {"type": "section", "title": "章"},
            {"type": "bullets", "title": "本文2", "bullets": ["い"]},
        ],
        "デッキ名",
    )
    numbers = []
    for slide in prs.slides:
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        numbers.append([t for t in texts if t.isdigit()])
    # Full-bleed slides carry no footer, and numbering skips them.
    assert numbers == [[], ["1"], [], ["2"]]


def test_deck_title_is_optional(tmp_path):
    prs = build_presentation([{"type": "bullets", "title": "見出し", "bullets": ["あ"]}])
    texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
    assert not any(t.isdigit() for t in texts)


@pytest.mark.parametrize(
    "spec, expected",
    [
        ({"type": "unknown"}, "unknown type"),
        ({"type": "bullets", "message": "1\n2\n3", "bullets": ["あ"]}, "up to 2 lines"),
        (
            {"type": "swimlane", "columns": [str(i) for i in range(11)], "lanes": [{"name": "a", "bars": []}]},
            "up to 10 columns",
        ),
        (
            {"type": "sequence", "actors": ["a", "b"], "messages": [{"from": 0, "to": 1, "label": "x"}] * 11},
            "up to 10 messages",
        ),
        (
            {"type": "sequence", "actors": ["a", "b"], "messages": [{"from": 0, "to": 5}]},
            "does not exist",
        ),
        (
            {"type": "swimlane", "columns": ["1", "2"], "lanes": [{"name": "a", "bars": [{"start": 1, "span": 5}]}]},
            "outside the 2 columns",
        ),
        ({"type": "agenda", "items": [f"項目{i}" for i in range(20)]}, "does not fit"),
        ({"type": "code", "code": "\n".join(f"line {i}" for i in range(30))}, "up to 18 lines"),
        ({"type": "columns", "columns": [{"title": "a", "items": []}]}, "2 to 4 columns"),
        ({"type": "table", "rows": [["a", "b"]]}, "at least one body row"),
        ({"type": "logictree", "root": {"label": "a"}, "direction": "sideways"}, 'direction must be "right" or "down"'),
        ({"type": "table", "rows": [["a", "b"], ["c"]]}, "same number of cells"),
        (
            {"type": "chart", "categories": ["1", "2"], "series": [{"name": "s", "values": [1]}]},
            "one value per category",
        ),
        (
            {"type": "chart", "chart_type": "donut", "categories": ["1"], "series": [{"name": "s", "values": [1]}]},
            "chart_type must be one of",
        ),
        (
            {"type": "quadrant", "points": [{"label": "x", "x": 1.4, "y": 0.5}]},
            "between 0 and 1",
        ),
    ],
)
def test_invalid_specs_raise_with_a_useful_message(spec, expected):
    with pytest.raises(ValueError, match=re.escape(expected)):
        build_presentation([spec])


def test_errors_name_the_slide_that_failed():
    with pytest.raises(ValueError, match=r"slide 2 \(table\)"):
        build_presentation(
            [
                {"type": "bullets", "title": "見出し", "bullets": ["あ"]},
                {"type": "table", "rows": [["a", "b"], ["c"]]},
            ]
        )


def test_empty_deck_is_rejected():
    with pytest.raises(ValueError, match="at least one slide"):
        build_presentation([])


def test_non_object_slide_is_rejected():
    with pytest.raises(ValueError, match="must be an object"):
        build_presentation(["タイトル"])


def test_deck_is_widescreen():
    prs = build_presentation([{"type": "title", "title": "表紙"}])
    assert prs.slide_width == int(SLIDE_W)
    assert prs.slide_height == int(SLIDE_H)


def test_east_asian_typeface_is_set(tmp_path):
    """Setting only the latin typeface leaves CJK glyphs on the theme font."""
    prs = build_presentation([{"type": "bullets", "title": "見出し", "bullets": ["日本語"]}])
    z, _ = parts(save(prs, tmp_path))
    root = etree.fromstring(z.read("ppt/slides/slide1.xml"))
    faces = {(e.tag.split("}")[1], e.get("typeface")) for e in root.iter() if e.tag.endswith(("}latin", "}ea"))}
    assert ("latin", "Yu Gothic") in faces
    assert ("ea", "Yu Gothic") in faces


def test_renderer_supplies_no_wording_of_its_own():
    """A deck must not pick up Japanese (or English) labels the caller did not ask
    for; GenU serves more than one language."""
    specs = [
        {"type": "chart", "title": "t", "categories": ["a"], "series": [{"name": "s", "values": [1]}], "insight": ["insight line"]},
        {"type": "sequence", "title": "t", "actors": ["a", "b"], "messages": [{"from": 0, "to": 1, "label": "call"}]},
        {"type": "swimlane", "title": "t", "columns": ["c1", "c2"], "lanes": [{"name": "lane", "bars": [{"start": 0, "span": 1, "label": "bar"}]}], "today": 1},
    ]
    # "deck" is the footer title and "1" the page number, both caller-driven.
    supplied = {"t", "a", "b", "c1", "c2", "s", "lane", "bar", "call", "insight line", "deck", "1"}
    for spec in specs:
        prs = build_presentation([spec], "deck")
        for shape in prs.slides[0].shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                assert text in supplied, f"{spec['type']} invented the label {text!r}"


def test_optional_labels_are_rendered_when_given():
    prs = build_presentation([SPECS["swimlane"]], "deck")
    texts = {s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame}
    assert "現在" in texts


def test_flow_branch_labels_default_to_english_and_are_overridable():
    default = build_presentation([SPECS["flow"]])
    texts = {s.text_frame.text for s in default.slides[0].shapes if s.has_text_frame}
    assert {"Yes", "No"} <= texts

    localised = build_presentation(
        [
            {
                **SPECS["flow"],
                "steps": [
                    {"type": "decision", "label": "条件", "no_branch": "別処理", "yes_label": "はい", "no_label": "いいえ"},
                ],
            }
        ]
    )
    texts = {s.text_frame.text for s in localised.slides[0].shapes if s.has_text_frame}
    assert {"はい", "いいえ"} <= texts


@pytest.mark.parametrize("name", sorted(SPECS))
def test_each_type_actually_draws_a_body(name):
    """A renderer that silently draws nothing leaves only the heading and footer,
    which every other assertion here would happily pass."""
    prs = build_presentation([SPECS[name]], "デッキ")
    chrome = {"見出し", "デッキ", "1"}
    body = [s for s in prs.slides[0].shapes if not (s.has_text_frame and s.text_frame.text.strip() in chrome)]
    assert len(body) >= 3, f"{name} drew nothing but the slide furniture"


def test_tree_directions_place_the_root_differently():
    right = build_presentation([SPECS["logictree"]], "デッキ")
    down = build_presentation([SPECS["logictree_down"]], "デッキ")

    def root_box(prs, label):
        return next(s for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text == label)

    def leaf_box(prs, label):
        return next(s for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text == label)

    # Going right the root is leftmost; going down it is topmost.
    assert root_box(right, "問い").left < leaf_box(right, "葉1").left
    assert root_box(down, "親").top < leaf_box(down, "孫1").top
    assert root_box(down, "親").left > leaf_box(down, "孫1").left


def _flow_steps(count):
    steps = [{"type": "start", "label": "開始"}]
    steps += [{"type": "process", "label": f"手順{i}"} for i in range(count - 2)]
    steps += [{"type": "end", "label": "終了"}]
    return steps


def test_a_long_flow_wraps_into_two_columns():
    """Wrapping beats shrinking the boxes until nobody can read them."""
    one = build_presentation([{"type": "flow", "title": "見出し", "steps": _flow_steps(5)}])
    many = build_presentation([{"type": "flow", "title": "見出し", "steps": _flow_steps(12)}])

    def node_columns(prs):
        lefts = [s.left for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text.startswith(("開始", "終了", "手順"))]
        return len(set(lefts))

    assert node_columns(one) == 1
    assert node_columns(many) == 2


def test_flow_columns_can_be_pinned():
    two = build_presentation([{"type": "flow", "title": "見出し", "columns": 2, "steps": _flow_steps(4)}])
    lefts = {s.left for s in two.slides[0].shapes if s.has_text_frame and s.text_frame.text.startswith(("開始", "終了", "手順"))}
    assert len(lefts) == 2

    with pytest.raises(ValueError, match="1 or 2 columns"):
        build_presentation([{"type": "flow", "columns": 3, "steps": _flow_steps(4)}])


def test_wrapping_a_flow_draws_nothing_higher_than_not_wrapping():
    """The connector between columns used to loop over the top of the slide,
    crossing the key message. The heading is identical either way, so the
    topmost shape must not move when the flow wraps."""
    common = {"type": "flow", "title": "見出し", "message": "一行"}
    one = build_presentation([{**common, "steps": _flow_steps(4)}])
    two = build_presentation([{**common, "steps": _flow_steps(9)}])

    def topmost(prs):
        return min(s.top for s in prs.slides[0].shapes)

    def column_count(prs):
        return len({s.left for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text.startswith(("開始", "終了", "手順"))})

    assert (column_count(one), column_count(two)) == (1, 2)
    assert topmost(two) >= topmost(one)


def test_a_flow_that_cannot_wrap_is_rejected():
    with pytest.raises(ValueError, match="does not fit"):
        build_presentation([{"type": "flow", "steps": _flow_steps(40)}])
